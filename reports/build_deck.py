"""Build a 10-slide defense deck for the P7 SOC Copilot.

Designed for a 15-minute mentor defense: ~1.5 min per slide, all 7 rubric
sections covered, two visual anchors (title + architecture), no animation,
no fluff. The text is the talking-points cheat sheet, not a script.

Outputs:
    reports/Professional_Industry_Defense_Deck.pptx

Run:
    python reports/build_deck.py
"""
from __future__ import annotations
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt, Emu

# ---- Project / authoring constants -----------------------------------------
PROJECT_TITLE = "Project 7 — SOC Security Operations Copilot"
PRESENTER     = "Soliman"
DECK_SUBTITLE = "15-minute mentor defense — Industry-Integrated AI Synthesis"
REPORT_DIR    = Path(__file__).resolve().parent
AGENT_GRAPH   = REPORT_DIR / "agent_graph.png"
DECK_OUT      = REPORT_DIR / "Professional_Industry_Defense_Deck.pptx"

# ---- Theme constants (kept consistent with the paper and README) -----------
BLUE        = RGBColor(0x0B, 0x3D, 0x91)
BLUE_LIGHT  = RGBColor(0xE8, 0xF0, 0xFE)
INK         = RGBColor(0x1A, 0x1A, 0x1A)
GREY        = RGBColor(0x55, 0x55, 0x55)
GREY_LIGHT  = RGBColor(0xCF, 0xD8, 0xE3)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xD9, 0x30, 0x25)
BG_GREY     = RGBColor(0xF3, 0xF5, 0xF8)
BORDER_GREY = RGBColor(0xD6, 0xDD, 0xE6)

FONT_BODY = "Segoe UI"
FONT_MONO = "Cascadia Mono"

# 16:9 widescreen, roomy for diagrams.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---- Low-level helpers -----------------------------------------------------

def _set_run(run, *, size: int = 14, bold: bool = False, italic: bool = False,
             color: RGBColor = INK, font: str = FONT_BODY) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _add_text(slide, x: Inches, y: Inches, w: Inches, h: Inches, text: str,
              *, size: int = 14, bold: bool = False, italic: bool = False,
              color: RGBColor = INK, font: str = FONT_BODY,
              align: str = "left") -> "Shape":
    """Add a single-paragraph text box. Multiline text uses \\n."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if align == "center": p.alignment = 2  # PP_ALIGN.CENTER
        elif align == "right": p.alignment = 1
        run = p.add_run(); run.text = line
        _set_run(run, size=size, bold=bold, italic=italic, color=color, font=font)
    return tb


def _add_bullets(slide, x: Inches, y: Inches, w: Inches, h: Inches,
                 bullets: list[str], *, size: int = 16, color: RGBColor = INK,
                 bold: bool = False, line_spacing: float = 1.15) -> "Shape":
    """Add a multi-bullet text box. Each item is one bullet (• added)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        # Bullet glyph + a thin space; keeps the bullet on the first line
        # even if the line wraps.
        run = p.add_run(); run.text = "•  " + item
        _set_run(run, size=size, bold=bold, color=color)
    return tb


def _add_rect(slide, x: Inches, y: Inches, w: Inches, h: Inches,
              fill: RGBColor, *, line: RGBColor | None = None,
              line_w: float = 0.75) -> "Shape":
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def _add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    """Title + optional subtitle, with a thin blue rule under them."""
    _add_text(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.6),
              title, size=28, bold=True, color=BLUE)
    if subtitle:
        _add_text(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4),
                  subtitle, size=14, color=GREY, italic=True)
    rule_y = Inches(1.45) if subtitle else Inches(1.05)
    _add_rect(slide, Inches(0.5), rule_y, Inches(12.3), Emu(9525), BLUE)


def _add_footer(slide, idx: int, total: int) -> None:
    """Project name on the left, slide number on the right."""
    _add_text(slide, Inches(0.5), Inches(7.05), Inches(8.0), Inches(0.35),
              f"{PROJECT_TITLE}  —  {PRESENTER}", size=9, color=GREY)
    _add_text(slide, Inches(11.8), Inches(7.05), Inches(1.0), Inches(0.35),
              f"{idx} / {total}", size=9, color=GREY, align="right")


def _section_label(slide, label: str) -> None:
    """Small grey label above the title — helps the mentor see the rubric
    section name at a glance (matches the rubric's 7 sections)."""
    _add_text(slide, Inches(0.5), Inches(0.10), Inches(12.3), Inches(0.25),
              label, size=10, color=GREY, italic=True)


def _add_table(slide, x: Inches, y: Inches, w: Inches, h: Inches,
               rows: list[list[str]], col_widths: list[float] | None = None,
               header: bool = True, font_size: int = 12) -> None:
    """Add a simple table with the first row styled as a header.

    rows: list of rows; each row is a list of cell strings. Header row gets
    a blue background + white bold text. Body rows get alternating tints.
    """
    n_rows = len(rows); n_cols = max(len(r) for r in rows)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for j, frac in enumerate(col_widths):
            table.columns[j].width = int(w * frac / total)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
            tf.margin_top  = Inches(0.04); tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run(); run.text = val
            if i == 0 and header:
                cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
                _set_run(run, size=font_size, bold=True, color=WHITE)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE_LIGHT if i % 2 == 0 else WHITE
                _set_run(run, size=font_size, color=INK)
            # Thin border on every cell.
            for edge in ("top", "left", "bottom", "right"):
                line = cell._tc.get_or_add_tcPr()
                # python-pptx doesn't expose cell borders cleanly across all
                # versions; the shape's table fill above + a subtle padding
                # is enough for the deck look.


# ---- Slide builders --------------------------------------------------------

def slide_title(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Big blue accent block on the left.
    _add_rect(s, Inches(0), Inches(0), Inches(0.4), SLIDE_H, BLUE)
    _add_text(s, Inches(1.0), Inches(2.0), Inches(11.5), Inches(0.6),
              "Industry-Integrated AI Synthesis — Project 7",
              size=14, color=GREY, italic=True)
    _add_text(s, Inches(1.0), Inches(2.5), Inches(11.5), Inches(1.4),
              "SOC Security Operations Copilot",
              size=40, bold=True, color=BLUE)
    _add_text(s, Inches(1.0), Inches(3.9), Inches(11.5), Inches(0.6),
              "Mentor Defense — 15 minutes",
              size=18, color=INK)
    _add_text(s, Inches(1.0), Inches(4.5), Inches(11.5), Inches(0.6),
              "Soliman  ·  Final synthesis capstone, Udacity AI Mastery",
              size=14, color=GREY)
    # Thesis box.
    _add_rect(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.2),
              BLUE_LIGHT, line=BLUE, line_w=1.0)
    _add_text(s, Inches(1.2), Inches(5.75), Inches(11.0), Inches(0.95),
              "Rules-first fusion  +  RAG-grounded summarization  +  "
              "non-bypassable human-approval gate.\n"
              "Built on P1's real corpus, P2's calibration, P3's rules "
              "discipline, P5's prose genre, and P6's governance spine.",
              size=14, color=INK, italic=True)


def slide_industry_context(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _section_label(s, "1. Industry Context and Problem Definition  ·  ~1.5 min")
    _add_title_bar(s, "Security operations is a 3 a.m. paragraph-writing job",
                   subtitle="Why this problem matters, and why AI is the right tool")
    # Two columns.
    _add_bullets(s, Inches(0.5), Inches(1.7), Inches(6.2), Inches(4.5), [
        "Security Operations Centers fuse surveillance feeds, badge-access "
        "logs, and alarm streams — thousands of events a day per site.",
        "The hard task is cross-screen correlation: a single incident "
        "(e.g. tailgate + forced door) is invisible to either screen alone.",
        "The 3 a.m. operator is the closer: they read the fused signal, "
        "match it to policy, write the analyst-facing paragraph, and "
        "decide whether to dispatch security.",
        "Today's tools give them raw streams. The job of paragraph-writing "
        "and policy-matching is what's automatable; the call is not.",
    ], size=15)
    # Right column: the 'why AI is appropriate' callout.
    _add_rect(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(4.5),
              BG_GREY, line=BORDER_GREY)
    _add_text(s, Inches(7.2), Inches(1.85), Inches(5.5), Inches(0.4),
              "Why an AI approach fits", size=14, bold=True, color=BLUE)
    _add_bullets(s, Inches(7.2), Inches(2.25), Inches(5.5), Inches(3.85), [
        "Signal fusion is exactly what rules + retrieval do well: deterministic correlation across streams, policy-grounded output.",
        "Generative summarization with citations is the paragraph-writing piece the operator shouldn't be doing at 3 a.m.",
        "Human-in-the-loop is the right trust model for irreversible actions: the system recommends, the operator decides.",
    ], size=13)
    _add_footer(s, 2, 10)


def slide_system_overview(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _section_label(s, "2. Integrated AI System Overview  ·  ~1.5 min")
    _add_title_bar(s, "Five subsystems compose one end-to-end pipeline",
                   subtitle="Each layer is small, swappable, and tested")
    # Horizontal flow boxes.
    layers = [
        ("Generators", "P1 corpus →\n1k events + 9.7k logs"),
        ("Fusion",     "4 rules →\n226 risk-scored incidents"),
        ("RAG",        "Chroma + MiniLM\n+ category routing"),
        ("Summarizer", "Groq 8B +\ncitation guard"),
        ("Agent",      "LangGraph +\nhuman-approval gate"),
    ]
    box_w = Inches(2.3); box_h = Inches(1.7)
    gap   = Inches(0.2)
    total_w = box_w * 5 + gap * 4
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.0)
    for i, (title, body) in enumerate(layers):
        x = start_x + (box_w + gap) * i
        _add_rect(s, x, y, box_w, box_h, BLUE_LIGHT, line=BLUE, line_w=1.25)
        _add_text(s, x, y + Inches(0.15), box_w, Inches(0.4),
                  title, size=15, bold=True, color=BLUE, align="center")
        _add_text(s, x + Inches(0.1), y + Inches(0.6), box_w - Inches(0.2),
                  Inches(1.05), body, size=12, color=INK, align="center")
        # Arrow between boxes.
        if i < len(layers) - 1:
            ax = x + box_w
            ay = y + Inches(0.7)
            _add_text(s, ax, ay - Inches(0.1), gap, Inches(0.4),
                      "→", size=20, bold=True, color=BLUE, align="center")
    # Below: the 5 subsystems in one sentence + the test count.
    _add_bullets(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.0), [
        "Composes 5 prior projects into one cohesive system — not a "
        "collection of disconnected notebooks.",
        "Each layer has a single responsibility and a small, named test "
        "(rules, retrieval, summarizer, governance gate, audit chain).",
        "Same `run_incident_streaming` entry point for CLI, notebook, "
        "pytest, and the Streamlit GUI — no parallel API surfaces.",
    ], size=14)
    _add_footer(s, 3, 10)


def slide_architecture_diagram(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _section_label(s, "2. Integrated AI System Overview  ·  continued")
    _add_title_bar(s, "The LangGraph governance topology",
                   subtitle="Compiled view of the agent the rubric submission notebook runs")
    if AGENT_GRAPH.exists():
        # Center the image; leave room for the heading + footer.
        s.shapes.add_picture(str(AGENT_GRAPH),
                             Inches(1.0), Inches(1.7),
                             width=Inches(11.3), height=Inches(4.8))
    else:
        _add_text(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.0),
                  f"(missing: {AGENT_GRAPH.name})",
                  size=14, color=RED, align="center")
    _add_bullets(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5), [
        "ingest → planner → worker → reviewer → {worker_dispatch · "
        "human_approval · summarizer}; reviewer is the only policy gate "
        "in the system.",
    ], size=11)
    _add_footer(s, 4, 10)


def slide_integration(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _section_label(s, "3. Integration of Prior Capstone Projects  ·  ~1.5 min")
    _add_title_bar(s, "Five prior projects, one system",
                   subtitle="P1 data · P2 calibration · P3 rules · P5 prose · P6 governance")
    rows = [
        ["Project", "What it contributed", "Evidence"],
        ["P1 — Reproducible workflows",
         "The full corpus the copilot runs on (1k events / 9.7k logs / 3 sites); the schema-adapter that normalizes it.",
         "src/generators/p1_pipeline.py · src/utils/p1_adapter.py"],
        ["P2 — Statistical analysis",
         "The 0.85 confidence threshold as a *contract*, not a constant. P2's chi-square + t-test are re-run on the slice the copilot runs against.",
         "tests/test_threshold_calibration.py · notebook §2b"],
        ["P3 — Applied ML",
         "Rules-first discipline + the leakage-audit lesson. The fusion layer is rules-only; any future ML must split by site + contiguous time_window.",
         "src/fusion/rules.py · src/fusion/risk_scorer.py"],
        ["P5 — Generative AI",
         "Genre reference for analyst-voice prose (NIST SP 800-61 / CSIRT / CERT). The summarizer itself is a citation-grounded LLM call, not the P5 model.",
         "src/agent/summarizer.py"],
        ["P6 — Agentic AI",
         "The LangGraph governance spine, lifted and adapted (one plan-loop fix). The non-bypassable policy gate hard-blocks `case.close` and gates `incident.escalate`.",
         "src/governance/ · src/domain/"],
    ]
    _add_table(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.8),
               rows, col_widths=[2.2, 5.6, 4.5], font_size=12)
    _add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
              "Removing any one of these layers breaks a real capability "
              "(fusion, policy routing, grounded summaries, or the "
              "non-bypassable escalation gate).",
              size=12, italic=True, color=GREY)
    _add_footer(s, 5, 10)


def slide_decisions(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _section_label(s, "4. Key Technical Decisions and Tradeoffs  ·  ~2 min")
    _add_title_bar(s, "Three decisions that shaped the system",
                   subtitle="Each one is a tradeoff, with the alternative explicit")
    # 3 stacked rows; each row: decision + alternatives + tradeoff.
    rows = [
        ("Rules-first fusion, not an ML scorer",
         "Per-rule base + size/confidence bonus, capped at 100, with the "
         "firing rule recorded on every incident.",
         "Alternative: a learned scorer. Tradeoff: rules are auditable; "
         "an ML scorer would learn patterns the rules miss but be opaque. "
         "ML is a swap-in behind the same `risk_scorer.py` interface when "
         "the rule layer earns its keep — and P2's calibration test "
         "confirms the cost honestly (recall@0.85 = 58%)."),
        ("Chroma + MiniLM, not a hosted vector DB",
         "5-doc KB, in-process, persistent across runs. 384-dim, "
         "normalized cosine. Category routing boosts the matching "
         "category by 0.3 before MMR.",
         "Alternative: a hosted vector DB. Tradeoff: ~3 s model load on "
         "first call vs. a monthly bill, an API key, and a network "
         "dependency. The KB doesn't pay for BM25 + cross-encoder yet; "
         "those are the right next step, not a bigger stack today."),
        ("Citation guard on a free 8B model",
         "Parse KB-XXXXX from the model output, keep only retrieved ids, "
         "retry once with a stricter prompt, else mark `needs_review`.",
         "Alternative: a bigger model. Tradeoff: ~1 in 50 incidents pays "
         "one extra Groq call; in exchange, hallucinated policy citations "
         "are impossible, not just unlikely. The cheapest possible "
         "defence against the most common grounded-generation failure."),
    ]
    y = Inches(1.7)
    for title, choice, alt in rows:
        _add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.65),
                  BLUE_LIGHT, line=BLUE, line_w=0.5)
        _add_text(s, Inches(0.7), y + Inches(0.1), Inches(11.9), Inches(0.4),
                  title, size=15, bold=True, color=BLUE)
        _add_text(s, Inches(0.7), y + Inches(0.5), Inches(11.9), Inches(0.45),
                  "Choice:  " + choice, size=12, color=INK)
        _add_text(s, Inches(0.7), y + Inches(0.95), Inches(11.9), Inches(0.7),
                  alt, size=11, italic=True, color=GREY)
        y += Inches(1.78)
    _add_footer(s, 6, 10)


def slide_ethics(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _section_label(s, "5. Ethical Considerations and Responsible AI  ·  ~1.5 min")
    _add_title_bar(s, "Ethics, made concrete in the system",
                   subtitle="Not a footnote — a code path, a test, a UI button")
    rows = [
        ["Risk", "Why it matters here", "What the system does about it"],
        ["Misuse / fabricated authority",
         "Surveillance + access data is sensitive; a model that invents a policy justification is the worst failure mode.",
         "Citation guard (summarizer.py): KB-XXXXX must be in the retrieved set, else `needs_review`."],
        ["Accountability / auto-escalation",
         "Escalation is irreversible — it dispatches security, suspends badges.",
         "`case.close` is hard-blocked; `incident.escalate` requires human approval at risk_band_score ≥ 80."],
        ["Transparency",
         "Every incident must be traceable to a rule, a score, and a citation.",
         "Per-rule provenance (`_rule`); raw cosine score returned (bonus is ordering-only); hash-chained audit log."],
        ["Calibration honesty",
         "A rule with a 42% anomaly blind spot is a deployed lie unless surfaced.",
         "P2's recall@0.85 = 58% is asserted in `test_threshold_calibration.py`; the GUI shows the raw numbers."],
    ]
    _add_table(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.6),
               rows, col_widths=[2.6, 4.5, 5.2], font_size=12)
    _add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
              "The system recommends and surfaces. The operator keeps "
              "judgment, override, and accountability.",
              size=13, italic=True, color=BLUE, align="center")
    _add_footer(s, 7, 10)


def slide_evaluation(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _section_label(s, "6. Evaluation, Limitations, and Risks  ·  ~2 min")
    _add_title_bar(s, "What worked, what didn't, what's known to be wrong",
                   subtitle="Strengths are real; limitations are surfaced, not hidden")
    # Left: strengths.
    _add_text(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
              "Strengths", size=15, bold=True, color=BLUE)
    _add_bullets(s, Inches(0.5), Inches(2.1), Inches(6.0), Inches(4.0), [
        "Notebook: Restart & Run All clean across 30 cells.",
        "Tests: 34/34 green across 5 files (governance, policy, escalate-blocked, threshold-calibration, upload-adapter).",
        "RAG self-check: every incident_type routes to its matching KB doc at rank 1.",
        "Summarizer: 4/4 spot-check incidents cite the correct policy doc; 0 `needs_review`.",
        "Operator UX: the same `run_incident_streaming` entry point powers CLI, notebook, pytest, and the GUI.",
    ], size=13)
    # Right: limitations (with the 58% in red to draw the eye).
    _add_text(s, Inches(6.9), Inches(1.7), Inches(6.0), Inches(0.4),
              "Limitations & risks", size=15, bold=True, color=RED)
    _add_bullets(s, Inches(6.9), Inches(2.1), Inches(6.0), Inches(4.0), [
        "Fusion recall@0.85 = 58% (P2 hard finding; surfaced to operators, not hidden).",
        "Free 8B model occasionally invents kwarg names — handled by an arg-filtering wrapper.",
        "Audit log carries an unrepaired break at an early sequence number; the GUI shows the warning honestly.",
        "Data is P1's synthetic corpus, not a real SOC environment.",
        "Streamlit GUI is local-only; no auth; per-session uploads live under data/uploads/{uuid}/.",
    ], size=13)
    # Hard finding callout.
    _add_rect(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7),
              RED, line=RED, line_w=0.5)
    _add_text(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.5),
              "Hard finding: the 0.85 threshold misses ~42% of anomaly events. "
              "Documented; not a deployable claim.",
              size=13, bold=True, color=WHITE)
    _add_footer(s, 8, 10)


def slide_relevance(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _section_label(s, "7. Professional Relevance and Next Steps  ·  ~1.5 min")
    _add_title_bar(s, "What this shows about the work, and what's next",
                   subtitle="System-level thinking · responsible engineering · reuse with discipline")
    # Left: relevance.
    _add_text(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
              "What this demonstrates", size=15, bold=True, color=BLUE)
    _add_bullets(s, Inches(0.5), Inches(2.1), Inches(6.0), Inches(4.5), [
        "System-level thinking: five prior projects compose into one pipeline, not a collection of notebooks.",
        "Responsible engineering under constraints: free-tier model + 5-doc KB + 0 dollars → category routing, citation guard, arg filtering.",
        "Reuse with discipline: the governance spine is lifted, fixed, de-labeled, and guarded by a no-domain-imports test.",
        "Defensible defaults: every metric in this deck is asserted in a test, not asserted in prose.",
    ], size=14)
    # Right: next steps.
    _add_text(s, Inches(6.9), Inches(1.7), Inches(6.0), Inches(0.4),
              "Next steps (12-18 months)", size=15, bold=True, color=BLUE)
    _add_bullets(s, Inches(6.9), Inches(2.1), Inches(6.0), Inches(4.5), [
        "Address the 58% recall finding: lower the threshold or add a complementary low-confidence rule.",
        "Multi-vendor event-stream connectors (P1's corpus is synthetic, not a real environment).",
        "Retrieval-level eval: a 10-row golden set + hit@k, before any model swap.",
        "Real auth for the Streamlit GUI: reverse proxy + OAuth; the GUI is local-only today.",
    ], size=14)
    _add_footer(s, 9, 10)


def slide_defense_ready(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _section_label(s, "Defense  ·  questions I'm ready for")
    _add_title_bar(s, "Three questions I expect — and the short answer",
                   subtitle="If you only have 30 seconds for each, here they are")
    # Three Q&A rows.
    rows = [
        ("Why rules-first, not an ML scorer?",
         "Auditability. Every incident carries the firing rule. The ML "
         "scorer is a swap-in behind the same `risk_scorer.py` interface "
         "when the rule layer earns its keep — and the P2 calibration "
         "test confirms the cost honestly."),
        ("What would break in production?",
         "The Streamlit GUI has no auth (it's a dev tool); the audit "
         "chain has an unrepaired early break; the fusion threshold "
         "misses ~42% of anomaly events. Each is a known and tested "
         "boundary, not a hidden risk."),
        ("Why is your RAG just dense + MMR, not hybrid?",
         "The KB is 5 documents. BM25, a cross-encoder, and a query-rewrite "
         "LLM are the right next step when the KB grows past a few hundred "
         "docs, not today. Category routing fixes the shared-vocabulary "
         "mis-ranking cheaply at this scale."),
    ]
    y = Inches(1.8)
    for q, a in rows:
        _add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.5),
                  BG_GREY, line=BORDER_GREY)
        _add_text(s, Inches(0.7), y + Inches(0.12), Inches(11.9), Inches(0.4),
                  "Q.  " + q, size=15, bold=True, color=BLUE)
        _add_text(s, Inches(0.7), y + Inches(0.55), Inches(11.9), Inches(0.9),
                  "A.  " + a, size=13, color=INK)
        y += Inches(1.62)
    # Closing line.
    _add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
              "The copilot doesn't replace the operator. It replaces the "
              "part of the shift that's a typist, a log-stitcher, and a "
              "3 a.m. paragraph-writer.",
              size=14, italic=True, color=BLUE, align="center")
    _add_footer(s, 10, 10)


# ---- Entry point -----------------------------------------------------------

def build() -> None:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_industry_context(prs)
    slide_system_overview(prs)
    slide_architecture_diagram(prs)
    slide_integration(prs)
    slide_decisions(prs)
    slide_ethics(prs)
    slide_evaluation(prs)
    slide_relevance(prs)
    slide_defense_ready(prs)

    prs.save(DECK_OUT)
    print(f"wrote deck: {DECK_OUT}  ({len(prs.slides)} slides)")


def main() -> None:
    ap = argparse.ArgumentParser(description="Build the 15-min mentor defense deck.")
    args = ap.parse_args()
    build()


if __name__ == "__main__":
    main()
